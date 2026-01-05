import pdfplumber
import requests
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


OLLAMA_URL = "http://localhost:11434/api/generate"
MODEL = "mistral"

PDF_FOLDER = "pdfs"
OUTPUT_PPT = "Cybersecurity_Summary_Report.pptx"


def extract_pdf_text(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def call_mistral(prompt):
    response = requests.post(
        OLLAMA_URL,
        json={
            "model": MODEL,
            "prompt": prompt,
            "stream": False
        }
    )
    return response.json()["response"]


def summarize_individual(text, filename):
    prompt = f"""
You are a cybersecurity analyst.

Summarize the following document ({filename}) focusing on:
- Threats
- Vulnerabilities
- Security incidents
- Recommendations

Use bullet points.
Document:
{text[:6000]}
"""
    return call_mistral(prompt)


def correlate_all(summaries):
    combined = "\n\n".join(summaries)

    prompt = f"""
You are a senior SOC analyst.

Given multiple cybersecurity report summaries below, correlate them and produce:
1. Common threats across clients
2. Repeated vulnerabilities
3. Patterns or attack trends
4. Overall risk level (Low / Medium / High)
5. Unified security recommendations

Summaries:
{combined}
"""
    return call_mistral(prompt)


def add_formatted_slide(prs, title_text, content_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    title = slide.shapes.title
    title.text = title_text
    tf = title.text_frame
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    # Content
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    lines = content_text.split("\n")

    for line in lines:
        line = line.strip()
        if not line:
            continue

        p = tf.add_paragraph()

        if line.lower().startswith(("-", "•")):
            p.text = line.lstrip("-• ").strip()
            p.font.size = Pt(14)
            p.level = 1
        else:
            p.text = line.replace("*", "")
            p.font.size = Pt(18)
            p.font.bold = True
            p.level = 0


def split_content_into_chunks(content, max_lines=12):
    lines = [line.strip() for line in content.split("\n") if line.strip()]
    chunks = []
    current = []

    for line in lines:
        current.append(line)
        if len(current) >= max_lines:
            chunks.append("\n".join(current))
            current = []

    if current:
        chunks.append("\n".join(current))

    return chunks


def main():
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Cybersecurity Assessment Summary"
    title_slide.placeholders[1].text = "AI-Generated Report using Mistral (Local LLM)"

    all_summaries = []

    for pdf_file in os.listdir(PDF_FOLDER):
        if pdf_file.endswith(".pdf"):
            print(f"Processing {pdf_file}...")
            text = extract_pdf_text(os.path.join(PDF_FOLDER, pdf_file))
            summary = summarize_individual(text, pdf_file)
            all_summaries.append(f"{pdf_file} Summary:\n{summary}")

    chunks = split_content_into_chunks(summary)

    for idx, chunk in enumerate(chunks, start=1):
        title = f"Summary – {pdf_file} ({idx}/{len(chunks)})"
        add_formatted_slide(prs, title, chunk)



    # Correlation slide
    correlated_summary = correlate_all(all_summaries)
    chunks = split_content_into_chunks(correlated_summary)

    for idx, chunk in enumerate(chunks, start=1):
        title = f"Correlated Cybersecurity Findings ({idx}/{len(chunks)})"
        add_formatted_slide(prs, title, chunk)



    prs.save(OUTPUT_PPT)
    print(f"\n✅ PPT generated successfully: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
